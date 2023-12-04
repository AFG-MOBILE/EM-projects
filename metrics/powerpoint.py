from pptx import Presentation
from copy import deepcopy
from pptx.dml.color import RGBColor
from pptx.util import Cm, Pt
from matplotlib import colors


def replace_hashtags_slide(path_presentation, data, path_new_presentation):
    replacements = {"=":"?"}
    # Carregar a apresentação do PowerPoint
    presentation = Presentation(path_presentation)
    for slide in presentation.slides:
        # Substituir marcadores nas tabelas
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            if '#' in paragraph.text:
                                text = paragraph.text
                                for chave, valor in data.items():
                                    text = text.replace(f'#{chave}#', str(valor))
                                paragraph.text = text
                                paragraph.font.size = Pt(11)
                                paragraph.font.color.rgb = RGBColor(116,35,133)
                                break

    # Salvar a apresentação com os marcadores preenchidos
    presentation.save(path_new_presentation)